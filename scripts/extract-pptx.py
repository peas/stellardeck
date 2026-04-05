#!/usr/bin/env python3
"""Extract content from PPTX files into Deckset-style markdown + webp images."""

import sys
import os
import io
import hashlib
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

ASSETS_DIR = "/Users/peas/presentations-paulo/assets"

def extract_pptx(pptx_path, output_md_path, image_prefix, orig_timestamp):
    """Extract slides from a PPTX into Deckset markdown."""
    print(f"Opening {pptx_path}...")
    prs = Presentation(pptx_path)

    slides_md = []
    image_count = 0
    hash_to_name = {}  # md5 -> image filename for dedup

    # Get original date for frontmatter
    orig_date = datetime.fromtimestamp(orig_timestamp).strftime("%Y-%m-%d")
    orig_date_full = datetime.fromtimestamp(orig_timestamp).strftime("%Y-%m-%d %H:%M:%S")

    total_slides = len(prs.slides)

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        if slide_num % 10 == 1 or slide_num == total_slides:
            print(f"  Slide {slide_num} de {total_slides}...")

        slide_texts = []
        slide_images = []
        slide_notes = ""

        # Extract notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text.strip():
                slide_notes = notes_frame.text.strip()

        # Process shapes
        for shape in slide.shapes:
            # Text
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        # Detect titles (larger font or placeholder type)
                        is_title = False
                        try:
                            if shape.placeholder_format is not None:
                                ph_type = shape.placeholder_format.type
                                # 1=TITLE, 15=CENTER_TITLE, 13=SUBTITLE
                                if ph_type in (1, 15):
                                    is_title = True
                        except (ValueError, AttributeError):
                            pass

                        if is_title and not text.startswith("#"):
                            slide_texts.append(f"# {text}")
                        else:
                            slide_texts.append(text)

            # Images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_blob = shape.image.blob
                    content_type = shape.image.content_type

                    # Hash to avoid duplicates
                    img_hash = hashlib.md5(image_blob).hexdigest()[:12]

                    if img_hash in hash_to_name:
                        # Reuse already-saved image
                        slide_images.append(hash_to_name[img_hash])
                        continue

                    image_count += 1
                    img_name = f"{image_prefix}-slide{slide_num:02d}-{image_count:03d}.webp"
                    img_path = os.path.join(ASSETS_DIR, img_name)

                    # Convert to webp
                    try:
                        img = Image.open(io.BytesIO(image_blob))
                        if img.mode in ('RGBA', 'LA', 'PA'):
                            img.save(img_path, 'WEBP', quality=85)
                        else:
                            img = img.convert('RGB')
                            img.save(img_path, 'WEBP', quality=85)
                        hash_to_name[img_hash] = img_name
                        slide_images.append(img_name)
                    except Exception as e:
                        print(f"    Warning: Could not convert image on slide {slide_num}: {e}")
                except Exception as e:
                    print(f"    Warning: Could not extract image from slide {slide_num}: {e}")

            # Group shapes (SmartArt, etc.)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for child in shape.shapes:
                        if child.has_text_frame:
                            for para in child.text_frame.paragraphs:
                                text = para.text.strip()
                                if text:
                                    slide_texts.append(text)
                except:
                    pass

            # Tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells]
                    slide_texts.append(" | ".join(row_texts))

        # Build slide markdown
        parts = []

        # Images as background (Deckset style)
        for img_name in slide_images:
            parts.append(f"![](../../assets/{img_name})")

        # Text
        for text in slide_texts:
            parts.append(text)

        # Notes
        if slide_notes:
            parts.append(f"\n^ {slide_notes}")

        if parts:
            slides_md.append("\n\n".join(parts))
        else:
            slides_md.append(f"<!-- Slide {slide_num}: empty or unsupported content -->")

    # Build final markdown
    header = f"""<!-- Original file: {os.path.basename(pptx_path)} -->
<!-- Original date: {orig_date_full} -->
<!-- Extracted: {datetime.now().strftime("%Y-%m-%d")} -->
<!-- {total_slides} slides, {image_count} unique images extracted -->

"""

    md_content = header + "\n\n---\n\n".join(slides_md)

    # Write markdown
    with open(output_md_path, 'w') as f:
        f.write(md_content)

    print(f"  Wrote {output_md_path} ({total_slides} slides, {image_count} images)")
    return total_slides, image_count


if __name__ == "__main__":
    base = "/Users/peas/presentations-paulo"

    files = [
        {
            "pptx": f"{base}/old/scuba-dev-devt/Alura_v2/MB3403_Alura_v2_fonte_DIN.pptx",
            "md": f"{base}/old/scuba-dev-devt/Alura_v2/MB3403_Alura_v2_fonte_DIN.md",
            "prefix": "alura-v2-din",
            "timestamp": 1599145152,  # Sep 3, 2020
        },
        {
            "pptx": f"{base}/old/scuba-dev-devt/Alura_v2/MB3403_Alura_v2_fonte_Arial.pptx",
            "md": f"{base}/old/scuba-dev-devt/Alura_v2/MB3403_Alura_v2_fonte_Arial.md",
            "prefix": "alura-v2-arial",
            "timestamp": 1599520205,  # Sep 7, 2020
        },
    ]

    for f in files:
        print(f"\n{'='*60}")
        print(f"Processing: {os.path.basename(f['pptx'])}")
        print(f"{'='*60}")
        slides, images = extract_pptx(f["pptx"], f["md"], f["prefix"], f["timestamp"])

        # Preserve original modification date
        touch_date = datetime.fromtimestamp(f["timestamp"]).strftime("%Y%m%d%H%M.%S")
        os.system(f"touch -t {touch_date} '{f['md']}'")
        print(f"  Set modification date to {touch_date}")

    print("\nDone!")
